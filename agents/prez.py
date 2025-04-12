import os
from typing import Optional
from dotenv import load_dotenv
import requests
from utils.prompts import pitch_deck_prompt
from pptx import Presentation
import json
from langchain_openai import ChatOpenAI

load_dotenv()


class PrezAgent():

  def __init__(self, marketing_docs, engineering_docs, endpoint=None):
    self.llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.9)
    self.marketing_docs = marketing_docs
    self.engineering_docs = engineering_docs

  def get_response(self):

    prompt = pitch_deck_prompt(self.marketing_docs, self.engineering_docs)
    print(f"Generating json for pitch deck...")
    response = self.llm.invoke(prompt).content
    print(f"JSON for pitch deck: {response}")
    return response

  def generate_pitch_deck(self):

    try:
      content_map = json.loads(self.get_response())
    except Exception as e:
      print(f"Error json loads for pitch deck: {e}")
      return None

    template_path = "utils/template.pptx"
    prs = Presentation(template_path)
    for i, slide in enumerate(prs.slides):
      print(f"\n--- Slide {i + 1} ---")

      for shape in slide.shapes:

        if shape.has_text_frame and not shape.name.startswith(
            "TextBox") and not shape.name.startswith("Freeform"):

          if shape.has_text_frame:
            # shape.text = content_map.get(shape.name, shape.text)
            tf = shape.text_frame

            para = tf.paragraphs[0]

            if para.runs:

              if shape.name in content_map:
                print(f"Updating text for shape: {shape.name}")
                para.runs[0].text = content_map[shape.name]

                for run in para.runs[1:]:
                  para._element.remove(run._r)

    prs.save("outputs/pitch_deck.pptx")
